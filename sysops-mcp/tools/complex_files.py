"""
Specialized file creation tools for complex formats (docx, xlsx, pptx, etc.). 
This module provides wrappers around specialized Python libraries to ensure reliable and structured content generation for non-plain-text files.

Dependencies required: python-docx, openpyxl, python-pptx, reportlab, pandas, PyYAML.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, List, Dict, Optional, Union
import csv
from datetime import datetime
from pydantic import BaseModel, Field
import traceback

# ---------------------------------------------------------------------------
# Schemas
# ---------------------------------------------------------------------------

class DocxConfig(BaseModel):
    """Configuration for DOCX files."""
    title: str = Field(..., description="Title of the document.")
    content: str = Field(..., description="Main body content.")
    sections: Optional[List[str]] = Field(default=None, description="Optional sections/subheadings.")


class XlsxConfig(BaseModel):
    """Configuration for XLSX files."""
    sheet_name: str = Field(default="Sheet1", description="Name of the active sheet.")
    data: List[List[Any]] = Field(..., description="Row by row data.")


class PptxConfig(BaseModel):
    """Configuration for PPTX files."""
    title: str = Field(..., description="Title of the presentation.")
    slides_content: List[str] = Field(..., description="Content for each slide.")


class PdfConfig(BaseModel):
    """Configuration for PDF files."""
    title: str = Field(..., description="Title of the PDF.")
    content: str = Field(..., description="Main body content.")


class CsvTsvConfig(BaseModel):
    """Configuration for CSV/TSV files."""
    headers: Optional[List[str]] = Field(default=None, description="Column headers.")
    data: List[List[Any]] = Field(..., description="Row by row data.")


class YamlConfig(BaseModel):
    """Configuration for YAML files."""
    data: Dict[str, Any] = Field(..., description="Dictionary data to serialize.")


class XmlConfig(BaseModel):
    """Configuration for XML files."""
    root_element: Dict[str, Any] = Field(..., description="Root element data to serialize.")


class ComplexFileRequest(BaseModel):
    """Unified schema for creating complex files."""
    file_type: str = Field(..., description="Type of file: docx, xlsx, pptx, pdf, csv, tsv, yaml, xml.")
    file_name: str = Field(..., description="Name of the file (e.g., 'report'). Extension is auto-added.")
    session_id: Optional[str] = Field(default=None, description="Optional session ID to segregate artifacts.")
    
    # Specific configs (one of these should be populated based on file_type)
    docx_config: Optional[DocxConfig] = None
    xlsx_config: Optional[XlsxConfig] = None
    pptx_config: Optional[PptxConfig] = None
    pdf_config: Optional[PdfConfig] = None
    csv_tsv_config: Optional[CsvTsvConfig] = None
    yaml_config: Optional[YamlConfig] = None
    xml_config: Optional[XmlConfig] = None


# ---------------------------------------------------------------------------
# Gateway functionalities
# ---------------------------------------------------------------------------

def _resolve_artifact_path(file_name: str, file_type: str, session_id: Optional[str] = None) -> Path:
    """
    Resolves the artifact path to a centralized directory segregated by session.
    Base directory is '~/mcp_artifacts'.
    """
    base_dir = Path.home() / "mcp_artifacts"

    if not session_id:
        session_id = datetime.now().strftime("%Y-%m-%d")

    session_dir = base_dir / session_id
    session_dir.mkdir(parents=True, exist_ok=True)

    if not file_name.lower().endswith(f".{file_type.lower()}"):
        file_name = f"{file_name}.{file_type.lower()}"

    return session_dir / file_name


def handle_complex_file_creation(request: ComplexFileRequest) -> dict[str, Any]:
    """
    Gateway function to handle the creation of complex files.
    """
    try:
        file_type = request.file_type.lower()
        path = _resolve_artifact_path(request.file_name, file_type, request.session_id)
        
        if file_type == "docx" and request.docx_config:
            return create_document(path, request.docx_config)
        elif file_type == "xlsx" and request.xlsx_config:
            return create_spreadsheet(path, request.xlsx_config)
        elif file_type == "pptx" and request.pptx_config:
            return create_presentation(path, request.pptx_config)
        elif file_type == "pdf" and request.pdf_config:
            return create_pdf(path, request.pdf_config)
        elif file_type in ["csv", "tsv"] and request.csv_tsv_config:
            return create_csv_tsv(path, request.csv_tsv_config, file_type)
        elif file_type == "yaml" and request.yaml_config:
            return create_yaml(path, request.yaml_config)
        elif file_type == "xml" and request.xml_config:
            return create_xml(path, request.xml_config)
        else:
            return {
                "success": False, 
                "error": f"Missing or mismatching configuration for file type: {file_type}"
            }
    except Exception as e:
        return {
            "success": False,
            "error": f"An error occurred during file creation: {str(e)}",
            "traceback": traceback.format_exc()
        }


def get_schema(file_type: Optional[str] = None) -> dict[str, Any]:
    """Returns the JSON schema for the specified file_type, or all if none provided."""

    schemas = """
Here is the complete expected schema for the complex file tool, starting with the base payload and 
breaking down into the individual configuration options for each file type:

### 1. Base Request Schema (`ComplexFileRequest`)
This is the root payload your LLM agent will need to provide to the `write_complex_file` tool.
```json
{
  "file_type": "string", // Required: "docx", "xlsx", "pptx", "pdf", "csv", "tsv", "yaml", or "xml"
  "file_name": "string", // Required: The name of the file (e.g., "financial_report")
  "session_id": "string", // Optional: Segregates artifacts. Defaults to current date (YYYY-MM-DD).
  
  // Provide ONLY ONE of the following configurations based on the 'file_type'
  "docx_config": { ... },
  "xlsx_config": { ... },
  "pptx_config": { ... },
  "pdf_config": { ... },
  "csv_tsv_config": { ... },
  "yaml_config": { ... },
  "xml_config": { ... }
}
```

---

### 2. Specific Configuration Schemas

#### DOCX Config (`docx_config`)
```json
{
  "title": "string", // Required: Title of the document (Heading 1)
  "content": "string", // Required: Main body paragraph content
  "sections": ["string", "string"] // Optional: Array of texts for sections/subheadings (Heading 2)
}
```

#### XLSX Config (`xlsx_config`)
```json
{
  "sheet_name": "string", // Optional: Defaults to "Sheet1"
  "data": [ // Required: 2D array representing row-by-row data
    ["Header 1", "Header 2", "Header 3"],
    [10, 20, 30],
    [40, 50, 60]
  ]
}
```

#### PPTX Config (`pptx_config`)
```json
{
  "title": "string", // Required: Title of the presentation (Title slide)
  "slides_content": [ // Required: Array of strings, where each string is the text content of a single slide
    "Content for Slide 1",
    "Content for Slide 2"
  ]
}
```

#### PDF Config (`pdf_config`)
```json
{
  "title": "string", // Required: Title of the PDF (Printed in bold at the top)
  "content": "string" // Required: Main body content string
}
```

#### CSV & TSV Config (`csv_tsv_config`)
```json
{
  "headers": ["Col 1", "Col 2"], // Optional: Array of column headers
  "data": [ // Required: 2D array representing row-by-row data
    ["Value 1A", "Value 2A"],
    ["Value 1B", "Value 2B"]
  ]
}
```

#### YAML Config (`yaml_config`)
```json
{
  "data": { // Required: A standard key-value JSON dictionary to be serialized into YAML
    "key1": "value1",
    "key2": ["list_item1", "list_item2"]
  }
}
```

#### XML Config (`xml_config`)
```json
{
  "root_element": { // Required: A dictionary mapping out the desired XML nodes starting inside the <root> element
    "parent_node": {
      "child_node": "value",
      "items": ["item1", "item2"]
    }
  }
}
```
    """

    return {"success": True, "schemas": schemas}


# ---------------------------------------------------------------------------
# Sub-functions (File Creators)
# ---------------------------------------------------------------------------

def create_document(path: Path, config: DocxConfig) -> dict[str, Any]:
    """Create DOCX document."""
    try:
        from docx import Document
    except ImportError:
        return {"success": False, "error": "The 'python-docx' library is not installed."}

    try:
        doc = Document()
        doc.add_heading(config.title, level=1)
        doc.add_paragraph(config.content)
        
        if config.sections:
            for section in config.sections:
                doc.add_heading(f'Section: {section}', level=2)
                doc.add_paragraph(section)

        doc.save(str(path))
        return {"success": True, "message": f"Successfully created DOCX file at {path}", "filepath": str(path)}
    except Exception as e:
        return {"success": False, "error": f"Failed to create document: {e}"}


def create_spreadsheet(path: Path, config: XlsxConfig) -> dict[str, Any]:
    """Create XLSX spreadsheet."""
    try:
        from openpyxl import Workbook
    except ImportError:
        return {"success": False, "error": "The 'openpyxl' library is not installed."}

    try:
        wb = Workbook()
        ws = wb.active
        ws.title = config.sheet_name

        for row_data in config.data:
            ws.append(row_data)

        wb.save(str(path))
        return {"success": True, "message": f"Successfully created XLSX file at {path}", "filepath": str(path)}
    except Exception as e:
        return {"success": False, "error": f"Failed to create spreadsheet: {e}"}


def create_presentation(path: Path, config: PptxConfig) -> dict[str, Any]:
    """Create PPTX presentation."""
    try:
        from pptx import Presentation
    except ImportError:
        return {"success": False, "error": "The 'python-pptx' library is not installed."}

    try:
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = config.title

        for i, content in enumerate(config.slides_content):
            blank_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.title.text = f"Slide {i+1}"
            slide.placeholders[1].text = str(content)

        prs.save(str(path))
        return {"success": True, "message": f"Successfully created PPTX presentation at {path}", "filepath": str(path)}
    except Exception as e:
        return {"success": False, "error": f"Failed to create presentation: {e}"}


def create_pdf(path: Path, config: PdfConfig) -> dict[str, Any]:
    """Create PDF document."""
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
    except ImportError:
        return {"success": False, "error": "The 'reportlab' library is not installed."}

    try:
        c = canvas.Canvas(str(path), pagesize=letter)
        textobject = c.beginText()
        textobject.setFont("Helvetica-Bold", 24)
        textobject.textLine(config.title)
        textobject.setY(700) 
        c.drawString(100, 750, config.title)

        textobject.setFont("Helvetica", 12)
        textobject.textLine(config.content)
        textobject.showWriter(c)

        c.save()
        return {"success": True, "message": f"Successfully created PDF file at {path}", "filepath": str(path)}
    except Exception as e:
        return {"success": False, "error": f"Failed to create PDF: {e}"}


def create_csv_tsv(path: Path, config: CsvTsvConfig, file_type: str) -> dict[str, Any]:
    """Create CSV or TSV file."""
    try:
        delimiter = '\t' if file_type == 'tsv' else ','
        with open(str(path), 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f, delimiter=delimiter)
            if config.headers:
                writer.writerow(config.headers)
            for row in config.data:
                writer.writerow(row)
        return {"success": True, "message": f"Successfully created {file_type.upper()} file at {path}", "filepath": str(path)}
    except Exception as e:
        return {"success": False, "error": f"Failed to create {file_type.upper()}: {e}"}


def create_yaml(path: Path, config: YamlConfig) -> dict[str, Any]:
    """Create YAML file."""
    try:
        import yaml
    except ImportError:
        return {"success": False, "error": "The 'PyYAML' library is not installed."}

    try:
        output_content = yaml.dump(config.data, sort_keys=False)
        with open(str(path), 'w', encoding='utf-8') as f:
            f.write(output_content)
        return {"success": True, "message": f"Successfully created YAML file at {path}", "filepath": str(path)}
    except Exception as e:
        return {"success": False, "error": f"Failed to write YAML data: {e}"}


def create_xml(path: Path, config: XmlConfig) -> dict[str, Any]:
    """Create XML file."""
    import xml.etree.ElementTree as ET
    try:
        root = ET.Element("root")
        
        def build_element(parent, data):
            if isinstance(data, dict):
                for key, value in data.items():
                    child = ET.SubElement(parent, key)
                    build_element(child, value)
            elif isinstance(data, list):
                for item in data:
                    child = ET.SubElement(parent, "item")
                    build_element(child, item)
            else:
                elem = ET.SubElement(parent, 'value')
                elem.text = str(data)

        build_element(root, config.root_element)
        
        tree = ET.ElementTree(root)
        ET.indent(tree, space="  ")
        tree.write(str(path), encoding='utf-8', xml_declaration=True)
        return {"success": True, "message": f"Successfully created XML file at {path}", "filepath": str(path)}
    except Exception as e:
        return {"success": False, "error": f"Failed to write XML data: {e}"}
